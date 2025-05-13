from pptx import Presentation

prs = Presentation()

layout = prs.slide_layouts[0]

slide = prs.slides.add_slide(layout)

title = slide.shapes[0]

subtitle = slide.shapes[1]

title.text = "Hello, World!"

subtitle.text = "python-pptx"

prs.save("./ppt/test.pptx")